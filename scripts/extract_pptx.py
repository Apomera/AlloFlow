from pptx import Presentation
import os

pptx_path = r'c:\Users\cabba\OneDrive\Desktop\UDL-Tool-Updated\AlloFlow_Architecture.pptx'
prs = Presentation(pptx_path)

output = []
for slide_num, slide in enumerate(prs.slides, 1):
    output.append(f"\n{'='*60}")
    output.append(f"SLIDE {slide_num}")
    output.append(f"{'='*60}")
    
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    output.append(f"  [{shape.shape_type}] {text}")
        
        if shape.has_table:
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                output.append(f"  [TABLE R{row_idx}] {' | '.join(row_text)}")
    
    if not any(s.startswith('  [') for s in output[-5:] if s.startswith('  ')):
        output.append("  (no text content)")

result = '\n'.join(output)
print(result[:8000])  # Print first 8000 chars

# Also save full text
with open(r'c:\Users\cabba\OneDrive\Desktop\UDL-Tool-Updated\scripts\pptx_text_dump.txt', 'w', encoding='utf-8') as f:
    f.write(result)

print(f"\nTotal slides: {len(prs.slides)}")
print(f"Full text saved to pptx_text_dump.txt")

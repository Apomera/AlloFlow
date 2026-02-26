from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

pptx_path = r'c:\Users\cabba\OneDrive\Desktop\UDL-Tool-Updated\AlloFlow_Architecture.pptx'
out_dir = r'c:\Users\cabba\OneDrive\Desktop\UDL-Tool-Updated\scripts\slide_images'
os.makedirs(out_dir, exist_ok=True)

prs = Presentation(pptx_path)

for slide_num, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            ext = img.content_type.split('/')[-1]
            if ext == 'jpeg':
                ext = 'jpg'
            filename = f'slide_{slide_num:02d}.{ext}'
            filepath = os.path.join(out_dir, filename)
            with open(filepath, 'wb') as f:
                f.write(img.blob)
            print(f'Extracted {filename} ({len(img.blob):,} bytes)')
            break

print(f'\nDone! {len(prs.slides)} slides extracted to {out_dir}')

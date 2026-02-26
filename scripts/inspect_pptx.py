from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = r'c:\Users\cabba\OneDrive\Desktop\UDL-Tool-Updated\AlloFlow_Architecture.pptx'
prs = Presentation(pptx_path)

print(f"Width: {prs.slide_width}, Height: {prs.slide_height}")
print(f"Total slides: {len(prs.slides)}")

for slide_num, slide in enumerate(prs.slides, 1):
    print(f"\n--- SLIDE {slide_num} ({len(slide.shapes)} shapes) ---")
    for shape in slide.shapes:
        info = f"  Shape: {shape.shape_type}, Name: '{shape.name}'"
        if hasattr(shape, 'width'):
            info += f", Size: {shape.width}x{shape.height}"
        if shape.has_text_frame:
            text = shape.text_frame.text[:100].strip()
            info += f", Text: '{text}'"
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or 'Picture' in shape.name or 'Image' in shape.name:
            info += " [IMAGE]"
            if hasattr(shape, 'image'):
                try:
                    info += f", Content-Type: {shape.image.content_type}, Size: {len(shape.image.blob):,} bytes"
                except:
                    pass
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            info += f" [GROUP with {len(shape.shapes)} sub-shapes]"
            for sub in shape.shapes:
                sub_info = f"    Sub: {sub.shape_type}, Name: '{sub.name}'"
                if sub.has_text_frame:
                    sub_info += f", Text: '{sub.text_frame.text[:80].strip()}'"
                print(sub_info)
        if shape.has_table:
            info += f" [TABLE {len(shape.table.rows)}x{len(shape.table.columns)}]"
        print(info)
